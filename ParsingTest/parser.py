"""
parser.py
---------
Document parser plug-in.

TODAY  — built-in fallback parser (PyMuPDF, python-docx, python-pptx, openpyxl).
TOMORROW — replace the import below with your friend's parser:

    from my_friends_parser import parse          # must match signature below

Expected signature:
    parse(file_path: str) -> str

The function receives the absolute path to a document file and must return
the extracted text as a plain Python string (no markdown, no formatting).

Supported types for the built-in fallback:
    .pdf  .docx  .doc  .txt  .pptx  .xlsx
"""
from __future__ import annotations

from pathlib import Path

from utils import clean_text, tokenize


# ── Swap this import to plug in a different parser ──────────────────────────
# from my_friends_parser import parse
# ────────────────────────────────────────────────────────────────────────────


def parse(file_path: str) -> str:
    """
    Extract plain text from a document file.
    Returns raw extracted text with no formatting or markdown.
    """
    import fitz              # PyMuPDF
    import docx as _docx    # python-docx
    from pptx import Presentation

    path = Path(file_path)
    ext  = path.suffix.lower()

    if ext == ".pdf":
        doc   = fitz.open(str(path))
        pages = [clean_text(page.get_text()) for page in doc]
        doc.close()
        return "\n".join(pages)

    if ext in (".docx", ".doc"):
        doc   = _docx.Document(str(path))
        lines = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t:
                        lines.append(t)
        return "\n".join(lines)

    if ext == ".txt":
        return path.read_text(encoding="utf-8")

    if ext == ".pptx":
        prs   = Presentation(str(path))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
        return "\n".join(lines)

    if ext == ".xlsx":
        import openpyxl
        wb    = openpyxl.load_workbook(str(path), data_only=True)
        lines = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_text = "  ".join(str(v) for v in row if v is not None)
                if row_text.strip():
                    lines.append(row_text)
        return "\n".join(lines)

    raise NotImplementedError(
        f"parse(): unsupported file type '{ext}'.\n"
        f"File attempted: {file_path}"
    )